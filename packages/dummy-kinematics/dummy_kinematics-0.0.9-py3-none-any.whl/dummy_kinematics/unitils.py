from decimal import Decimal
from pptx import Presentation
from colorama import Fore, Style, init
import pandas as pd
import numpy as np
import copy
from scipy import signal
from typing import NoReturn, Tuple
from pptx.presentation import Presentation as _PresentationObj
from collections.abc import Iterable
from collections import Counter
from scipy.integrate import cumulative_trapezoid
from typing import Sequence
from numpy.typing import NDArray

init(autoreset=True)  # 彩打完成后默认重置恢复默认样式


def pptx_identify(prs_path: str, output_file: bool = True) -> NoReturn:
    """识别pptx模板占位符,打印输出占位符情况，默认输出确认文件。
    prs_path: 字符串输入需要确认的模板文件，如"template.pptx"
    output_file: bool,是否输出文件,默认为输出文件True。
    """
    prs: _PresentationObj = Presentation(prs_path)
    for idx, layout in enumerate(prs.slide_layouts):
        print(Style.BRIGHT +
              f"正在确认第{idx}个模板-{layout.name}, placeholder输出格式:idx-type-name")
        slide = prs.slides.add_slide(layout)
        for holder in slide.placeholders:
            phf = holder.placeholder_format
            print(f'   {phf.idx} -{phf.type} -{holder.name}')
            holder.text = f'模板{idx}-{layout.name}-ph{phf.idx} - {phf.type} -{holder.name}'
    if output_file:
        info = """按照模板序号-模板名-占位符idx-占位符类型 - 占位符名称写进该文件!
已经保存为confirm_pptx.pptx, 请查阅.
        """
        prs.save("confirm_pptx.pptx")
        print(Fore.RED + info)


def lazyproperty(func):
    name = '_lazy_' + func.__name__

    @property
    def lazy(self):
        if hasattr(self, name):
            return getattr(self, name)
        else:
            value = func(self)
            setattr(self, name, value)
            return value

    return lazy


def is_positive(value) -> bool | Iterable[bool]:
    """value : 可以为简单数值，或迭代类型
    """
    if not isinstance(value, Iterable):
        return value > 0
    else:
        return tuple(i > 0 for i in value)


def is_most_positive(value) -> bool:
    """判断一个值(简单值或可迭代值)是否大多数为正数
    Parameters
    ----------
    value :如1,[1,2]...

    Returns
    -------
    bool
    """
    positive_bool = is_positive(value)
    if not isinstance(positive_bool, Iterable):
        return positive_bool
    count = Counter(positive_bool)
    # 返回最多出现的bool值
    return count.most_common(1)[0][0]


def g_to_v(acc_g: NDArray, dx: float,
           initial_kph: float,
           transformed: bool = False) -> NDArray:
    """输出单位为m/s
    Parameters
    ----------
    acc_g : NDArray
        _description_
    initial_kph : float
        积分初始速度kph
    dx : float
        采样间隔ms, 一般是0.1
    transformed: bool
        是否自动转换方向

    Returns
    -------
    NDArray : ND
       速度,单位m/s
    """
    if not isinstance(acc_g, np.ndarray):
        raise TypeError("acc_g must be np.ndarray")
    acc_g = acc_g * 9.8  # ms2和g换成标准单位
    if acc_g.ndim == 1:
        if transformed and (not is_most_positive(acc_g)):  # 正碰， X要求acc为正值
            acc_g = -acc_g
        initial_speed = initial_kph / 3.6  # 转换程m/s
        v = initial_speed - \
            cumulative_trapezoid(acc_g, dx=0.001 * dx,
                                 initial=0, axis=0)  # v的单位是m/s
        return v

    if acc_g.ndim == 2:
        row_num, col_num = acc_g.shape
        for col in range(col_num):
            acc = acc_g[:, col]
            if not is_most_positive(acc):  # 要求acc为正值
                acc_g[:col] = -acc
            initial_speed = initial_kph / 3.6  # 转换程m/s
            v = initial_speed - \
                cumulative_trapezoid(acc_g, dx=0.001 * dx,
                                     initial=0, axis=0)  # v的单位是m/s
            return v
    else:
        raise ValueError("input data must ndim <=2")


def v_to_s(v_values: NDArray, dx: float, initial_mm: float = 0) -> NDArray:
    """根据速度m/s,计算位移mm

    Parameters
    ----------
    v_values : Iterable,m/s 为单位
    dx : float
        采样间隔ms, 如0.1
    initial_mm : float, optional
        初始位移mm, by default 0

    Returns
    -------
    Iterable
        位移mm
    """
    if not isinstance(v_values, np.ndarray):
        raise TypeError("input data must be np.ndarray")
    res = initial_mm + cumulative_trapezoid(v_values, dx=dx, initial=0, axis=0)
    return res


def g_to_s(acc_g: NDArray, dx: float,
           initial_kph: float,
           initial_mm: float = 0,
           transformed: bool = False) -> NDArray:
    """简化g_to_v, v_to_s

    Parameters
    ----------
    acc_g : Iterable
        _description_
    dx : float
        采样间隔ms,一般为0.1
    initial_kph : float
        _description_
    initial_mm : float, optional
        _description_, by default 0
    transformed: bool
        是否自动转换方向

    Returns
    -------
    Iterable
        _description_
    """
    v = g_to_v(acc_g=acc_g,
               dx=dx,
               initial_kph=initial_kph,
               transformed=transformed
               )
    s = v_to_s(v, dx=dx,
               initial_mm=initial_mm)
    return s


def simple_integrated_once(y_values: Iterable, dx: float, initial_value: float = 0):
    """不涉及单位转换的简单一次积分，给定采样数据间隔（固定）.特别注意dx如果不涉及单位转换，0.1ms如以
    国际单位制为0.1*10^-3s
    Parameters
    ----------
    y_values : Iterable
        _description_
    dx : float
        _description_
    initial_value : float, optional
        _description_, by default 0

    Returns
    -------
    _type_
        _description_
    """
    res = initial_value + cumulative_trapezoid(y_values, dx=dx, initial=0)
    return res


def deg_v_to_deg(degps: Iterable, dx: float, initial_deg: float = 0) -> Iterable:
    """角速度到角位移

    Parameters
    ----------
    degps : Iterable
        被积角速度数据 deg/s 
    dx : float
       采样间隔ms数
    initial_deg : float, optional
        _description_, by default 0

    Returns
    -------
    Iterable
        _description_
    """
    res = simple_integrated_once(degps, dx * 0.001, initial_deg)
    return res


def convert_to_most_positive_2darray(value: NDArray) -> NDArray:
    """转换成每列为负曲线列,一般用于acc_g列数据统一方向为负值
    Parameters
    ----------
    value : NDArray
    Returns
    """
    if value.ndim != 2:
        raise ValueError("input must be 2D np.array")
    row_num, col_num = value.shape
    for i in range(col_num):
        col = value[:, i]
        if not is_most_positive(col):
            value[:, i] = -col
    return value


def cfc_filter(cfc: int, data: NDArray, dt_ms: float, axis: int = 0):
    """
    Create filter using CFC based on SAE J211
    https://docs.scipy.org/doc/scipy/reference/generated/scipy.signal.filtfilt.html
    """
    dt = dt_ms * 0.001
    wd = 2 * 3.14159 * cfc * 2.0775
    wa = (np.sin(wd * dt / 2)) / (np.cos(wd * dt / 2))
    aO = (wa ** 2) / (1 + wa * (2 ** 0.5) + wa ** 2)
    a1 = 2 * aO
    a2 = aO
    b1 = -2 * ((wa ** 2) - 1) / (1 + wa * (2 ** 0.5) + wa ** 2)
    b2 = (-1 + wa * (2 ** 0.5) - wa ** 2) / (1 + wa * (2 ** 0.5) + wa ** 2)
    A = [1, -b1, -b2]
    B = [aO, a1, a2]
    filtered_data = signal.filtfilt(B, A, data, axis=axis)
    return filtered_data


def index_number(li: Iterable, target: float) -> Tuple[int, float]:
    """返回从序列中找最接近给定值的数值及索引，返回（索引，目标值） 

    :param _type_ li: 目标序列
    :param _type_ target: 目标值
    :return Tuple[int,float]: （索引，结果值）
    """

    res = min(enumerate(li), key=lambda tup: abs(tup[1] - target))
    return res


def cal_olc(acc_g: NDArray,
            v0_kph: float,
            start_ms: float = 0,
            end_ms: float = 200,
            step_ms: float = 0.1,
            cfc180_filter: bool = False,
            detailed_return: bool = False,
            ) -> float | dict:
    """根据波形计算OLC

    :param NDArray acc_g: 输入acc_g, 需要为一维NDARRAY
    :param float v0_kph: 初始速度 kph
    :param float start_ms: 给定波形的初始采用时间ms, defaults to 0
    :param float end_ms: 采样截至时间ms, defaults to 200
    :param float step_ms: 采样率时间ms, defaults to 0.1
    :param bool cfc180_filter: 如果是未滤波的波形,则开启滤波CFC180, defaults to False
    :param bool detailed_return: 详细模式返回，默认关闭
    :return _type_: olc 值,单位为g. 详细模式返回字典(olc,t1,v1,t2,v2,summary),t_ms, v_m/s字典
    """
    new = pd.DataFrame({"input_acc_g": acc_g})
    times = np.arange(start_ms, end_ms + step_ms, step_ms)
    new["times"] = times
    # CFC180滤波 如果有
    if cfc180_filter:
        filtered = cfc_filter(CFC=180, data=acc_g, dt_ms=step_ms)

    new["x"] = filtered if cfc180_filter else acc_g

    new["v_v0"] = v0_kph / 3.6
    v_v0 = new["v_v0"]
    v_mpdb = g_to_v(acc_g=new["x"].values,
                    initial_kph=v0_kph,
                    dx=step_ms,
                    transformed=True)
    new["v_mpdb"] = v_mpdb

    s_v0 = v_to_s(v_values=new["v_v0"].values, dx=step_ms)
    new["s_v0"] = s_v0
    s_mpdb = g_to_s(acc_g=new["x"].values,
                    dx=step_ms,
                    initial_kph=v0_kph,
                    transformed=True
                    )
    new["s_mpdb"] = s_mpdb

    sr_s_v0_s_mpdb = s_v0 - s_mpdb
    new["sr_s_v0_s_mpdb"] = sr_s_v0_s_mpdb

    # 求t1, 未约束假人运动65mm
    index, val = index_number(sr_s_v0_s_mpdb, 65)
    t1_index = copy.deepcopy(index)

    t1 = times[index]  # ms
    v1 = v0_kph / 3.6  # m/s
    print(f"t1: {t1:.2f}ms v1:{v1:.2f} m/s")

    new["gt_t1"] = new["times"] > t1

    # 因为安全带作用减小的假人位移，小于等于t1时候为0， 大于t1时候为三角形面积
    # 记作 sr_s_belt_effect, 这个不是真实的，这个只是映射关系， OLC不同，这个值就不同 不是随时间变化的

    def make_cal_belt_s(t1):
        v0 = v0_kph / 3.6

        def cal_belte_s(row: Iterable):
            if row['gt_t1']:
                res = 0.5 * (v0 - row["v_mpdb"]) * (row['times'] - t1)
                return res
            return 0

        return cal_belte_s

    new['sr_s_belt_effect'] = new.apply(make_cal_belt_s(t1), axis=1)
    sr_s_belt_effect = new["sr_s_belt_effect"].values
    # occupan 相对台车的位移
    sr_occupant = s_v0 - s_mpdb - sr_s_belt_effect
    new["sr_occupant"] = sr_occupant

    # v2,t2 时刻就是在虚拟假人安全带约束下继续向前235mm ， 假人总共300mm位移的时刻
    # 求t2,v2
    index, val = index_number(sr_occupant, 300)
    t2_index = copy.deepcopy(index)
    t2 = times[index]
    v2 = v_mpdb[index]
    print(f"t2:{t2:.2f}ms v2:{v2:.2f}m/s")
    olc = 1000 * (v1 - v2) / (t2 - t1) / 9.81  # 换算成单位g
    print(f"olc为{olc:.3f}g")

    # 计算虚拟假人的速度曲线，为了出图

    def make_occupant_v(row):
        if row['times'] < t1:
            return v1
        if row['times'] > t2:
            return v2
        t = row['times']
        k = (v1 - v2) / (t2 - t1)
        return v1 - (t - t1) * k

    new["v_occupant"] = new.apply(make_occupant_v, axis=1)
    v_occupant = new["v_occupant"].values

    # 计算虚拟假人的位移 ，为了出图
    new["s_occupant"] = v_to_s(new["v_occupant"].values, dx=0.1)
    s_occupant = new["s_occupant"].values

    # 交互环境下可以出图，以下仅在交互环境有图
    # new.plot.line(x="times", y=["v_mpdb", "v_v0", "v_occupant"])
    # new.plot.line(x="times", y=["s_v0", "s_occupant", "s_mpdb"])

    if not detailed_return:
        return olc
    # 生成简要说明
    summary = f"olc={olc:.2f}g,v1@t1={v1:.2f}m/s@{t1:.2f}ms,v2@t2={v2:.2f}m/s@{t2:.2f}ms"
    res = dict(olc=olc,
               t1=t1,
               v1=v1,
               t2=t2,
               v2=v2,
               v_v0=v_v0,
               v_mpdb=v_mpdb,
               v_occupant=v_occupant,
               s_mpdb=s_mpdb,
               s_v0=s_v0,
               s_occupant=s_occupant,
               summary=summary,
               df=new
               )
    return res


def s_to_v(s: Iterable,
           dx_ms: float = 0.1,
           initial_kph: float = None) -> NDArray:
    """s 单位为mm, 输出速度单位为m/s

    :param Iterable s: _description_
    :param float dx_ms: _description_, defaults to 0.1
    :param float initial_kph: _description_, defaults to None
    :return NDArray: _description_
    """
    intial = initial_kph / 3.6 if initial_kph is not None else 0
    res = [intial]
    for index, val in enumerate(s):
        if index == 0:
            continue
        else:
            diff = val - s[index - 1]  # mm
            vt = diff / dx_ms
            res.append(vt)
    return np.array(res)


def liner_insert(x: float, high_score_index: tuple,
                 low_score_index: tuple) -> float:
    """线性插值函数，如胸部高性能（4分，35mm) , 低性能（0分，60mm）,求40.6mm分
    则liner_insert(40.6,(4,35),(0,60))

    :param float x: 插值数
    :param tuple high_score_index: 高性能分数及对应限值
    :param tuple low_score_index:低性能分数及对应限值
    :return float: 分数
    """
    high_score, high_index = high_score_index
    low_score, low_index = low_score_index
    if x < high_index:
        return high_score
    if x > low_index:
        return low_score
    k = (high_score - low_score) / (high_index - low_index)
    return high_score + k * (x - high_index)


def cal_rttf(acc_g: NDArray,
             speed_kph: float,
             start_ms: float = 0,
             end_ms: float = 200,
             step_ms: float = 0.1) -> tuple[float, float]:
    """根据三英寸，五英寸-30ms 计算RTTF"""
    df = pd.DataFrame({"acc_g": acc_g})
    df["dummy_g"] = 0
    df['svel_abs'] = g_to_s(acc_g, step_ms, speed_kph, 0, True)
    df['sd_abs'] = g_to_s(df['dummy_g'].values, step_ms, speed_kph, 0)
    df['sd_relative'] = df['sd_abs'] - df['svel_abs']
    timeseries = np.arange(start_ms, end_ms + step_ms, step_ms)
    inch3 = 76
    inch5 = 127
    index, number = index_number(df['sd_relative'], inch3)
    rttf_min = timeseries[index] - 30
    rttf_min = round(rttf_min, 2)
    index, number = index_number(df['sd_relative'], inch5)
    rttf_max = timeseries[index] - 30
    rttf_max = round(rttf_max, 2)
    return rttf_min, rttf_max




def increase_index_number(li: Sequence, target: float):
    res = (0, li[0])
    gap = abs(res[1] - target)
    for idx, value in enumerate(li):
        if idx == 0:
            continue
        elif (abs(value - target) < gap) and ((value - li[idx - 1]) > 0):
            gap = abs(value - target)
            res = (idx, value)
    return res


def _init_seg_points(x: NDArray, *segs) -> list[tuple[int, float]]:
    """
    :param x:
    :param segs: 分割点
    :return: 返回[(index, number)...]
    """

    x_min = min(x)
    x_max = max(x)
    inputs = [x_min, *segs, x_max]
    inputs.sort()
    res = [increase_index_number(x, target) for target in inputs]
    return res


def get_seg_mean(seg_to_last_intg: list[float], seg_points_values: list[float]) -> list[float]:
    assert len(seg_to_last_intg) == len(seg_points_values), "inputs list should be ths same len"
    li_len = len(seg_points_values)
    res = []
    for i in range(li_len):
        try:
            res.append(seg_to_last_intg[i] / (seg_points_values[i] - seg_points_values[i - 1]))
        except IndexError:
            res.append(0)
    return res


def make_output_xy(x: NDArray, index_series: list[int], seg_to_last_mean: list[float]) -> tuple:
    """

    :param x:
    :param index_series:
    :param seg_to_last_mean:
    :return: 返回处理后的x,y , 均为 NDArray
    """
    assert len(index_series) == len(seg_to_last_mean), "input list should be the same length"
    out_x = x[:max(index_series) + 1]
    out_y = np.empty(max(index_series) + 1)
    index_mean = list(zip(index_series, seg_to_last_mean))

    def match_y(index: int):
        for idx1, mean in index_mean:
            if index <= idx1:
                return mean

    for idx, value in enumerate(out_y):
        out_y[idx] = match_y(idx)

    return out_x, out_y


def _get_seg_intg(x: NDArray, y: NDArray, start_index: int, end_index: int) -> float:
    from scipy.integrate import trapezoid
    filtered_y = y[start_index:end_index + 1]
    filtered_x = x[start_index:end_index + 1]
    return trapezoid(y=filtered_y, x=filtered_x)


def get_seg_to_last_intg(x: NDArray, y: NDArray, index_series: list[int]) -> list[float]:
    seg_to_last_intg = []
    for i in range(len(index_series)):
        if i == 0:
            seg_to_last_intg.append(0)

        else:
            seg_to_last_intg.append(_get_seg_intg(x, y, start_index=index_series[i - 1], end_index=index_series[i]))

    return seg_to_last_intg


def make_seg_mean(x: NDArray, y: NDArray, *segs) -> dict:
    """

    :param x:
    :param y:
    :param segs:
    :return: out_x=out_x, out_y=out_y, result_dict=result_dict
    """
    result_dict = {}
    segs_tuples = _init_seg_points(x, *segs)
    result_dict["seg_pt_values"] = [tup[1] for tup in segs_tuples]
    result_dict["index"] = [tup[0] for tup in segs_tuples]
    result_dict["seg_to_last_intg"] = get_seg_to_last_intg(x, y, result_dict["index"])
    result_dict["seg_to_last_mean"] = get_seg_mean(result_dict["seg_to_last_intg"], result_dict["seg_pt_values"])

    out_x, out_y = make_output_xy(x, result_dict["index"], result_dict["seg_to_last_mean"])
    return dict(out_x=out_x, out_y=out_y, result_dict=result_dict)






