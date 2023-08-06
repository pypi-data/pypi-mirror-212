
from pptx import Presentation
from colorama import Fore, Style, init
from typing import NoReturn, List, Self
from pptx.shapes.placeholder import PicturePlaceholder
from pptx.presentation import Presentation as _PresentationObj
from datetime import datetime

from .chart_element import Chartinfo
import os

init(autoreset=True)  # 彩打完成后默认重置恢复默认样式


class PTT_Page_inserter:

    def __init__(self):
        self._path = os.path.join(os.path.dirname(__file__),
                                  r"default\default_template.pptx")  # 模板文件地址，内置
        self.prs: _PresentationObj = Presentation(self._path)
        self._internal_pptx_identify()

    def _internal_pptx_identify(self, output_file: bool = False) -> NoReturn:
        """识别pptx模板占位符,打印输出占位符情况，默认不输出确认文件。
        prs_path: 字符串输入需要确认的模板文件,若为None,则为内置PPT模板路径
        output_file: bool,是否输出文件,默认为不输出文件。
        """

        prs = Presentation(self._path)  # 重新生成prs ,避免对正式报告产生影响
        for idx, layout in enumerate(prs.slide_layouts):
            print(Style.BRIGHT +
                  f"正在确认第{idx}个模板-{layout.name}, placeholder输出格式:idx-type-name")
            slide = prs.slides.add_slide(layout)
            for holder in slide.placeholders:
                phf = holder.placeholder_format
                print(f'   {phf.idx} -{phf.type} -{holder.name}')
                holder.text = f'模板{idx}-{layout.name}-ph{phf.idx} - {phf.type} -{holder.name}'
        if output_file:
            info = """按照模板序号-模板名-占位符idx-占位符类型 - 占位符名称写进该文件!
    已经保存为confirm_pptx.pptx, 请查阅.
            """
            prs.save("confirm_pptx.pptx")
            print(Fore.RED+info)

    def insert_charts_in2_pages(self, page_title: str, charts_info: List[Chartinfo],
                                conclusion: str = None) -> Self:
        """根据图表数据信息生成PPT页

        Parameters
        ----------
        page_title : str
            页面标题名称
        charts_info : List[Chartinfo]
            图表数据信息列表

        Returns
        -------
        Presentation
        返回PPT文件对象
        """

        layout = self.prs.slide_layouts[1]

        for idx, chart_info in enumerate(charts_info):
            if idx % 2 == 0:  # 偶数图，从0开始，需要增加一页PPT，且使用左侧的13号placeholder
                slide = self.prs.slides.add_slide(layout)
                slide.shapes.title.text = page_title
                if conclusion is not None:
                    conclusion_ph = slide.placeholders[15]
                    conclusion_ph.text = conclusion
                holder = slide.placeholders[13]
                comments = chart_info.chart_info_comment()
                slide.notes_slide.notes_text_frame.text = comments

            else:  # 从0开始数， 奇数图用右侧14号placeholder
                holder = slide.placeholders[14]
                comments = comments+chart_info.chart_info_comment()
                # 保存本页备注 最大值最小值信息,
                slide.notes_slide.notes_text_frame.text = comments

            chart_info.inserted(holder)

        return self

    def insert_charts_with_seat_info(self, page_title: str,
                                     charts_info: List[Chartinfo],
                                     seat_info: str = None,
                                     conclusion: str = None) -> Self:
        """根据图表数据信息生成PPT页
        seat_info:str 座椅位置及假人信息字符串，如 "DR: AM50"

        Returns
        -------
        Presentation
        返回PPT文件对象
        """

        layout = self.prs.slide_layouts[3]

        for idx, chart_info in enumerate(charts_info):
            if idx % 2 == 0:  # 偶数图，从0开始，需要增加一页PPT，且使用左侧的13号placeholder
                slide = self.prs.slides.add_slide(layout)
                if seat_info is not None:
                    seat_ph = slide.placeholders[16]
                    seat_ph.text = seat_info
                if conclusion is not None:
                    conclusion_ph = slide.placeholders[15]
                    conclusion_ph.text = conclusion
                slide.shapes.title.text = page_title
                holder = slide.placeholders[13]
                comments = chart_info.chart_info_comment()
                slide.notes_slide.notes_text_frame.text = comments

            else:  # 从0开始数， 奇数图用右侧14号placeholder
                holder = slide.placeholders[14]
                comments = comments+"\n"+chart_info.chart_info_comment()
                # 保存本页备注 曲线最大值最小值信息,
                slide.notes_slide.notes_text_frame.text = comments

            chart_info.inserted(holder)
        return self

    def insert_picture_pages(self, page_title: str, pic_paths: list[str]) -> Self:
        """图片插入到ppt生成PPT页

        Parameters
        ----------
        page_title : str
            页面标题名称
        pic_paths : list[str]
            图片地址路径集合
        Returns
        -------
        Presentation
        返回PPT文件对象
        """

        # self.internal_pptx_identify()

        layout = self.prs.slide_layouts[2]

        for idx, pic_path in enumerate(pic_paths):

            if idx % 2 == 0:

                slide = self.prs.slides.add_slide(layout)
                holder: PicturePlaceholder = slide.placeholders[16]
                slide.shapes.title.text = page_title
            else:
                holder: PicturePlaceholder = slide.placeholders[17]
            holder.insert_picture(pic_path)

        return self

    def insert_cover_page(self, title: str) -> Self:

        layout = self.prs.slide_layouts[0]

        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        datetime_holder = slide.placeholders[13]  # "powerd and datetime"
        datetime_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S ")
        datetime_holder.text = datetime_str
        return self
