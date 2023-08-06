from colorama import init
from typing import  NoReturn, Iterable, Literal
from pptx.chart.data import XyChartData
from pptx.chart.chart import Chart
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Pt
from typing import Self
from pptx.shapes.placeholder import ChartPlaceholder
from .data_converter import CSV_Data_Converter
from .unitils import cal_olc


init(autoreset=True)  # 彩打完成后默认重置恢复默认样式


class Seriesinfo:
    def __init__(self, name: str, x_data: Iterable, y_data: Iterable):
        self.name = name
        self.x_data = x_data
        self.y_data = y_data

    @property
    def max_value_str(self):
        x, y = max(zip(self.x_data, self.y_data),
                   key=lambda cop: float(cop[1]))

        res = "{:.2f}@{:.2f}ms".format(y, x)

        return res

    @property
    def min_value_str(self):
        x, y = min(zip(self.x_data, self.y_data),
                   key=lambda cop: float(cop[1]))
        res = "{:.2f}@{:.2f}ms".format(y, x)

        return res

    def series_info_comment(self) -> str:
        return f'  Series: {self.name} ,max={self.max_value_str}, min={self.min_value_str}'


class Chartinfo:
    """保存chartdata相关数据
    """

    def __init__(self, title: str, many_series_info: list[Seriesinfo],
                 x_axis_title: str = None,
                 y_axis_title: str = None) -> None:

        self.title = title
        self.many_series_info = many_series_info
        self.x_axis_title = x_axis_title
        self.y_axis_title = y_axis_title

    def making_XYchart_Data(self) -> XyChartData:
        """返回chartdata对象,可供placeholder insert_chart时使用

        Returns
        -------
        XyChartData
            返回实例对象
        """
        xychartdata = XyChartData()
        for series_info in self.many_series_info:
            series = xychartdata.add_series(series_info.name)
            for x, y in zip(series_info.x_data, series_info.y_data):
                series.add_data_point(x, y)
        return xychartdata

    def chart_info_comment(self) -> str:
        series_summary = ""
        for series in self.many_series_info:
            series_summary = series_summary+series.series_info_comment()+"\n"
        return f'Chart: {self.title} \n {series_summary}'

    def _style_chart_2(self, chart: Chart,
                       ) -> NoReturn:
        """设定chart样式,并添加标题,添加轴标题等,适合两张图
        """
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.has_title = True
        chart.chart_title.text_frame.text = self.title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.size = Pt(10)
        if self.x_axis_title:
            chart.category_axis.axis_title.text_frame.text = self.x_axis_title
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(
                10)
        if self.y_axis_title:
            chart.value_axis.axis_title.text_frame.text = self.y_axis_title
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(
                10)

    def _style_chart_4(self, chart: Chart,
                       ) -> NoReturn:
        """设定chart样式,并添加标题,添加轴标题等,适合4张图
        """
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)
        chart.has_title = True
        chart.chart_title.text_frame.text = self.title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)
        if self.x_axis_title:
            chart.category_axis.axis_title.text_frame.text = self.x_axis_title
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(
                8)
        if self.y_axis_title:
            chart.value_axis.axis_title.text_frame.text = self.y_axis_title
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(
                8)

    def inserted(self, c_ph: ChartPlaceholder, small: bool = False) -> NoReturn:
        chart: Chart = c_ph.insert_chart(
            XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS, self.making_XYchart_Data()).chart

        # 以下为设置图例样式和图标题
        if not small:
            self._style_chart_2(chart)
        if small:
            self._style_chart_4(chart)

    @classmethod
    def created_by_iso(cls, iso_code: str, data_converter_maps: dict,
                       transformed: bool = False) -> Self | None:
        """根据iso_code , 和试验数据(即data_coverter_maps),生成图chartinfo数据。
        若生成成功则返回chartinfo, 否则None

        Parameters
        ----------
        iso_code : str
        data_converter_maps : dict
        Returns
        -------
        Self | None
        """
        series_list = []
        completed = False  # 是否完成图
        for k, data_converter in data_converter_maps.items():
            data_converter: CSV_Data_Converter = data_converter
            y_data = data_converter.find_data_only_by_iso(iso_code,
                                                          transformed=transformed)
            if y_data is not None:  # 找到数据才添加到序列
                ser = Seriesinfo(
                    name=k, x_data=data_converter.time_series, y_data=y_data)
                series_list.append(ser)
                if not completed:  # 仅当图增加第一个序列的时候需要保存图名和y坐标轴单位名
                    chart_title = data_converter.find_chart_title_by_iso(
                        iso_code)
                    unit = data_converter.find_unit_by_iso(iso_code)
                    completed = True
        if completed:  # 具备生成图的条件
            chart_info = cls(title=chart_title,
                             many_series_info=series_list,
                             x_axis_title="time [ms]",
                             y_axis_title=unit)
        return chart_info if completed else None

    @classmethod
    def body_GTVTST_chartinfo(cls,
                              direction: Literal["x", "y", "z"], value_type: Literal["g", "v", "s"],
                              data_coverter_maps: dict[str, CSV_Data_Converter]) -> Self:
        chart_title = f"body {value_type.upper()}T-{direction.upper()} direction "
        series_list = []
        x_axis_title = "time [ms]"
        match value_type:
            case "g":
                y_axis_title = "accelerate [g]"
            case "v":
                y_axis_title = "velocity [m/s]"
            case "s":
                y_axis_title = "displacemnet [mm]"
            case _:
                raise ValueError(
                    'value_type parament must be Literal["g", "v", "s"]')

        for name, data_converter in data_coverter_maps.items():
            v = data_converter.body_value_getter(
                direction=direction, value_type=value_type)
            series = Seriesinfo(name=name, x_data=data_converter.time_series,
                                y_data=v)
            series_list.append(series)
        return cls(title=chart_title, many_series_info=series_list,
                   x_axis_title=x_axis_title,
                   y_axis_title=y_axis_title)

    @classmethod
    def body_GS_chartinfo(cls,
                          direction: Literal["x", "y", "z"],
                          data_coverter_maps: dict[str, CSV_Data_Converter]) -> Self:
        chart_title = f"body GS-{direction.upper()} direction"
        series_list = []
        x_axis_title = "displacement [mm]"
        y_axis_title = "accelerate [g]"
        for name, data_converter in data_coverter_maps.items():
            s = data_converter.body_value_getter(
                direction=direction, value_type="s")
            g = data_converter.body_value_getter(
                direction=direction, value_type="g")
            series = Seriesinfo(name=name, x_data=s,
                                y_data=g)
            series_list.append(series)
        return cls(title=chart_title, many_series_info=series_list,
                   x_axis_title=x_axis_title,
                   y_axis_title=y_axis_title)

    @classmethod
    def body_XY_ss_chartinfo(cls,
                          data_coverter_maps: dict[str, CSV_Data_Converter]) -> Self:
        """
        SOL 常用SS图， ACU XY SS图
        """
        chart_title = "body ACU S-S XY abs_displacement"
        series_list = []
        x_axis_title = "displacement X [mm]"
        y_axis_title = "displacement Y [mm]"
        for name, data_converter in data_coverter_maps.items():
            sx = data_converter.body_value_getter(
                direction="x", value_type="s")
            sy = data_converter.body_value_getter(
                direction="y", value_type="s")
            series = Seriesinfo(name=name, x_data=sx,
                                y_data=sy)
            series_list.append(series)
        return cls(title=chart_title, many_series_info=series_list,
                   x_axis_title=x_axis_title,
                   y_axis_title=y_axis_title)

    @classmethod
    def OLC_VTST_chartinfos_dict(cls, data_converter: CSV_Data_Converter,
                                 specify_isocode: str = None) -> dict |None:
        """ OLC 图生成，如果指定isocode ，则使用该ISO波形，否则使用自带车体波形X向
        返回字典， 字典为{vt_chart: ..., st_chart:...,summary: OLC summary}
        :param dict[str, CSV_Data_Converter] data_coverter_maps: 
        :param str specify_isocode:  defaults to None
        :return : 
        """
        name = data_converter.name
        # 一个datacoverter 就是一张图
        acc_g = data_converter.body_x_g
        v0_kph = data_converter.speed_kph
        start_ms = data_converter.start
        end_ms = data_converter.end
        step_ms = data_converter.step
        times = data_converter.time_series
        chart_title = f"{name} body OLC VT chart"
        if specify_isocode is not None:
            chart_title = f"{name} {specify_isocode} OLC VT chart"
            acc_g = data_converter.find_data_only_by_iso(specify_isocode)
            if acc_g is None:
                return None
        olc_dict = cal_olc(acc_g=acc_g, v0_kph=v0_kph, start_ms=start_ms, end_ms=end_ms, step_ms=step_ms,
                           detailed_return=True)
        # 制作结论
        summary = olc_dict["summary"]
        # 制作VT图
        body_vt_series = Seriesinfo(name="body_v" if not specify_isocode else f"{specify_isocode}_v",
                                    x_data=times,
                                    y_data=olc_dict["v_mpdb"])
        v0_vt_series = Seriesinfo(
            name="v0", x_data=times, y_data=olc_dict["v_v0"])
        v_occupant_series = Seriesinfo(
            name="v_occupant", x_data=times, y_data=olc_dict["v_occupant"])
        vt_chart = cls(title=chart_title,
                       many_series_info=[body_vt_series,
                                         v0_vt_series, v_occupant_series],
                       x_axis_title="time [ms]",
                       y_axis_title="velocity [m/s]")
        # 制作ST图
        chart_title = f"{name} body OLC ST chart"
        if specify_isocode is not None:
            chart_title = f"{name} {specify_isocode} OLC ST chart"
        s_v0_series = Seriesinfo(name="v0_s",
                                 x_data=times,
                                 y_data=olc_dict["s_v0"]
                                 )
        s_body_series = Seriesinfo(name="body_s" if not specify_isocode else f"{specify_isocode}_s",
                                   x_data=times,
                                   y_data=olc_dict['s_mpdb'])
        s_occupant_series = Seriesinfo(
            name="s_occupant", x_data=times, y_data=olc_dict["s_occupant"])
        st_chart = cls(title=chart_title,
                       many_series_info=[s_body_series,
                                         s_v0_series, s_occupant_series],
                       x_axis_title="time [ms]",
                       y_axis_title="displacement [mm]")

        one = {'vt_chart': vt_chart,
               "st_chart": st_chart,
               "summary": summary}
        return one

    @classmethod
    def created_by_meta(cls, meta_dict: dict) -> Self:
        """mate_dict:
        {
            title:"chart A",
            x_axis_title:"x_aixs"
            y_axis_title:"y_aix",
            series_meta:[
                {   name: "series 1",
                    x_data:[1,2],
                    y_data:[3,4]
                },...],
        }
        """
        series_meta_list = meta_dict.pop("series_meta")
        many_series_list = []
        for series_meta in series_meta_list:
            seri = Seriesinfo(**series_meta)
            many_series_list.append(seri)
        return cls(many_series_info=many_series_list, **meta_dict)

    @staticmethod
    def relative_vt_chart_meta(iso_code: str, data_converter_maps: dict[str, CSV_Data_Converter],
                               initial_kph_mode: bool) -> dict:
        direction = iso_code[-2]
        title = "relative V-T " + direction + "-direction"
        x_axis_title = "time [ms]"
        y_axis_title = "relative velocity [m/s]"
        series_meta = []
        one = {}
        for serise_name, data_converter in data_converter_maps.items():  # 添加多次试验的数据系列
            body_v = data_converter.body_value_getter(
                direction=direction, value_type="v")
            one["name"] = serise_name
            one["x_data"] = data_converter.time_series
            one["y_data"] = data_converter.g_to_v_by_iso(
                iso_code, initial_kph_mode)-body_v
            series_meta.append(one)
            one = {}  # 必须清空 否则返回的全是最后一次迭代的dict
        chart_meta = {
            "title": title,
            "x_axis_title": x_axis_title,
            "y_axis_title": y_axis_title,
            "series_meta": series_meta
        }
        return chart_meta

    @staticmethod
    def abs_vt_chart_meta(iso_code: str, data_converter_maps: dict[str, CSV_Data_Converter],
                          initial_kph_mode: bool) -> dict:
        direction = iso_code[-2]
        title = "abs V-T " + direction + "-direction"
        x_axis_title = "time [ms]"
        y_axis_title = "abs velocity [m/s]"
        series_meta = []
        one = {}
        for serise_name, data_converter in data_converter_maps.items():  # 添加多次试验的数据系列

            one["name"] = serise_name
            one["x_data"] = data_converter.time_series
            one["y_data"] = data_converter.g_to_v_by_iso(
                iso_code, initial_kph_mode)
            series_meta.append(one)
            one = {}  # 必须清空 否则返回的全是最后一次迭代的dict
        chart_meta = {
            "title": title,
            "x_axis_title": x_axis_title,
            "y_axis_title": y_axis_title,
            "series_meta": series_meta
        }
        return chart_meta

    @staticmethod
    def abs_st_chart_meta(iso_code: str, data_converter_maps: dict[str, CSV_Data_Converter],
                          initial_kph_mode: bool) -> dict:
        direction = iso_code[-2]
        title = "abs S-T " + direction + "-direction"
        x_axis_title = "time [ms]"
        y_axis_title = "abs displacement [mm]"
        series_meta = []
        one = {}
        for serise_name, data_converter in data_converter_maps.items():  # 添加多次试验的数据系列

            one["name"] = serise_name
            one["x_data"] = data_converter.time_series
            one["y_data"] = data_converter.g_to_s_by_iso(
                iso_code, initial_kph_mode)
            series_meta.append(one)
            one = {}  # 必须清空 否则返回的全是最后一次迭代的dict
        chart_meta = {
            "title": title,
            "x_axis_title": x_axis_title,
            "y_axis_title": y_axis_title,
            "series_meta": series_meta
        }
        return chart_meta

    @staticmethod
    def relative_st_chart_meta(iso_code: str, data_converter_maps: dict[str, CSV_Data_Converter],
                               initial_kph_mode: bool) -> dict:
        direction = iso_code[-2]
        title = "relative S-T " + direction + "-direction"
        x_axis_title = "time [ms]"
        y_axis_title = "relative displacement"
        series_meta = []
        one = {}
        for serise_name, data_converter in data_converter_maps.items():  # 添加多次试验的数据系列
            body_s = data_converter.body_value_getter(
                direction=direction, value_type="s")
            one["name"] = serise_name
            one["x_data"] = data_converter.time_series
            one["y_data"] = data_converter.g_to_s_by_iso(
                iso_code, initial_kph_mode)-body_s
            series_meta.append(one)
            one = {}  # 必须清空 否则返回的全是最后一次迭代的dict
        chart_meta = {
            "title": title,
            "x_axis_title": x_axis_title,
            "y_axis_title": y_axis_title,
            "series_meta": series_meta
        }
        return chart_meta

    @staticmethod
    def gs_relative_chart_meta(iso_code: str, data_converter_maps: dict[str, CSV_Data_Converter],
                               initial_kph_mode: bool) -> dict:
        direction = iso_code[-2]
        title = "G-S(relative) " + direction + "-direction"
        x_axis_title = "relative displacement[mm]"
        y_axis_title = "acc [g]"
        series_meta = []
        one = {}
        for serise_name, data_converter in data_converter_maps.items():  # 添加多次试验的数据系列
            body_s = data_converter.body_value_getter(
                direction=direction, value_type="s")
            one["name"] = serise_name
            one["x_data"] = data_converter.g_to_s_by_iso(
                iso_code, initial_kph_mode)-body_s
            one["y_data"] = data_converter.find_data_only_by_iso(iso_code,
                                                                 transformed=True)
            series_meta.append(one)
            one = {}  # 必须清空 否则返回的全是最后一次迭代的dict
        chart_meta = {
            "title": title,
            "x_axis_title": x_axis_title,
            "y_axis_title": y_axis_title,
            "series_meta": series_meta
        }
        return chart_meta

    @staticmethod
    def gs_abs_chart_meta(iso_code: str, data_converter_maps: dict[str, CSV_Data_Converter],
                          initial_kph_mode: bool) -> dict:
        direction = iso_code[-2]
        title = "G-S(absolute) " + direction + "-direction"
        x_axis_title = "relative displacement[mm]"
        y_axis_title = "acc [g]"
        series_meta = []
        one = {}
        for serise_name, data_converter in data_converter_maps.items():  # 添加多次试验的数据系列

            one["name"] = serise_name
            one["x_data"] = data_converter.g_to_s_by_iso(
                iso_code, initial_kph_mode)
            one["y_data"] = data_converter.find_data_only_by_iso(iso_code,
                                                                 transformed=True)
            series_meta.append(one)
            one = {}  # 必须清空 否则返回的全是最后一次迭代的dict
        chart_meta = {
            "title": title,
            "x_axis_title": x_axis_title,
            "y_axis_title": y_axis_title,
            "series_meta": series_meta
        }
        return chart_meta
